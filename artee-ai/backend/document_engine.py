import os
import fitz  # PyMuPDF
import pandas as pd
import openpyxl
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors

class DocumentIntelligenceAgent:
    """
    Classifies documents and extracts key structured metadata from PDFs, images, and text.
    """
    
    def classify_document(self, file_path: str) -> dict:
        """
        Scans file properties and textual content to classify invoice vs utility bill vs PO.
        """
        if not os.path.exists(file_path):
            return {"type": "unknown", "confidence": 0.0, "error": "File does not exist."}
            
        ext = os.path.splitext(file_path)[1].lower()
        
        # Extract sample text (supporting PDF text extraction)
        sample_text = ""
        if ext == ".pdf":
            try:
                doc = fitz.open(file_path)
                for page in doc[:3]: # check first 3 pages
                    sample_text += page.get_text()
            except Exception as e:
                print(f"[DocumentIntelligenceAgent] Failed reading PDF: {e}")
        elif ext in [".txt", ".csv", ".json"]:
            try:
                with open(file_path, "r", encoding="utf-8") as f:
                    sample_text = f.read(1000)
            except Exception:
                pass

        sample_lower = sample_text.lower()
        
        # Classification rules engine
        if "invoice" in sample_lower or "tax invoice" in sample_lower:
            return {"type": "invoice", "confidence": 0.95}
        elif "utility bill" in sample_lower or "electric" in sample_lower or "power bill" in sample_lower:
            return {"type": "utility_bill", "confidence": 0.90}
        elif "purchase order" in sample_lower or "po number" in sample_lower:
            return {"type": "purchase_order", "confidence": 0.92}
        elif "packing list" in sample_lower:
            return {"type": "packing_list", "confidence": 0.88}
            
        return {"type": "document", "confidence": 0.50}

    def extract_pdf_data(self, file_path: str) -> dict:
        """
        Parses invoice text layouts to extract totals and line items dynamically.
        """
        result = {"vendor": "Unknown", "amount": 0.0, "tax": 0.0, "items": []}
        try:
            doc = fitz.open(file_path)
            full_text = ""
            for page in doc:
                full_text += page.get_text()
            
            # Simple keyword extractor logic
            lines = full_text.split("\n")
            for line in lines:
                line_lower = line.lower()
                if "total" in line_lower or "net amount" in line_lower:
                    parts = line.split()
                    for p in parts:
                        try:
                            # Strip currency symbols and parse float
                            val = float(p.replace("$", "").replace(",", "").replace("₹", ""))
                            if val > result["amount"]:
                                result["amount"] = val
                        except ValueError:
                            pass
                elif "tax" in line_lower:
                    parts = line.split()
                    for p in parts:
                        try:
                            val = float(p.replace("$", "").replace(",", ""))
                            if val > result["tax"]:
                                result["tax"] = val
                        except ValueError:
                            pass
        except Exception as e:
            print(f"[DocumentIntelligenceAgent] Extraction error: {e}")
            
        return result


class ExcelAgent:
    """
    Cleans data sheets, removes duplicates, and performs price lists vendor comparisons.
    """

    def clean_and_merge_sheets(self, file_paths: list, output_path: str) -> bool:
        """
        Merges sheets from multiple price files, formatting amounts and removing duplicates.
        """
        dataframes = []
        for path in file_paths:
            if os.path.exists(path):
                try:
                    df = pd.read_excel(path)
                    # Normalize columns
                    df.columns = [str(c).strip().lower() for c in df.columns]
                    dataframes.append(df)
                except Exception as e:
                    print(f"[ExcelAgent] Failed reading sheet {path}: {e}")

        if not dataframes:
            return False

        try:
            combined = pd.concat(dataframes, ignore_index=True)
            # Remove exact duplicate rows
            combined.drop_duplicates(inplace=True)
            
            # Clean empty rows
            combined.dropna(how="all", inplace=True)
            
            combined.to_excel(output_path, index=False)
            return True
        except Exception as e:
            print(f"[ExcelAgent] Merge failed: {e}")
            return False

    def compare_vendor_prices(self, file_path: str) -> list:
        """
        Scans a combined spreadsheet and recommends vendors sorted by lowest unit cost.
        """
        recommendations = []
        try:
            df = pd.read_excel(file_path)
            # Find item column, price column, vendor column
            cols = list(df.columns)
            item_col = next((c for c in cols if "item" in c or "description" in c), cols[0])
            price_col = next((c for c in cols if "price" in c or "cost" in c or "val" in c), cols[1])
            vendor_col = next((c for c in cols if "vendor" in c or "supplier" in c or "store" in c), cols[2] if len(cols) > 2 else cols[0])

            # Group items and find min prices
            for name, group in df.groupby(item_col):
                sorted_group = group.sort_values(by=price_col)
                cheapest = sorted_group.iloc[0]
                recommendations.append({
                    "item": str(name),
                    "cheapest_price": float(cheapest[price_col]),
                    "vendor": str(cheapest[vendor_col])
                })
        except Exception as e:
            print(f"[ExcelAgent] Price list comparison failed: {e}")

        return recommendations


class OfficeAutomationAgent:
    """
    Generates professional styled reports using python-docx, python-pptx, and ReportLab.
    """

    def generate_word_report(self, output_path: str, title: str, sections: dict):
        """
        Creates a custom formatted DOCX report.
        """
        doc = Document()
        doc.add_heading(title, 0)
        
        for subtitle, content in sections.items():
            doc.add_heading(subtitle, level=1)
            doc.add_paragraph(content)
            
        doc.save(output_path)
        return True

    def generate_powerpoint_presentation(self, output_path: str, title: str, slides_content: list):
        """
        Creates a PowerPoint PPTX file representing slides sequence.
        """
        prs = Presentation()
        
        # Add title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title_box = slide.shapes.title
        subtitle_box = slide.placeholders[1]
        title_box.text = title
        subtitle_box.text = "Artee AI Operations Report"
        
        # Add content slides
        for slide_data in slides_content:
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            
            title_shape.text = slide_data.get("heading", "Section Title")
            tf = body_shape.text_frame
            tf.text = slide_data.get("body", "")
            
        prs.save(output_path)
        return True

    def generate_pdf_report(self, output_path: str, title: str, data: list):
        """
        Generates a premium styled PDF document using ReportLab SimpleDocTemplate.
        """
        doc = SimpleDocTemplate(output_path, pagesize=letter, rightMargin=40, leftMargin=40, topMargin=40, bottomMargin=40)
        styles = getSampleStyleSheet()
        
        # Custom styles definition
        title_style = ParagraphStyle(
            name='TitleStyle',
            fontName='Helvetica-Bold',
            fontSize=22,
            leading=26,
            textColor=colors.HexColor('#0f172a'),
            spaceAfter=20
        )
        body_style = ParagraphStyle(
            name='BodyStyle',
            fontName='Helvetica',
            fontSize=11,
            leading=15,
            textColor=colors.HexColor('#334155'),
            spaceAfter=15
        )
        
        story = []
        story.append(Paragraph(title, title_style))
        story.append(Spacer(1, 10))
        story.append(Paragraph("This procurement report summarizes vendor price checks, and lists recommended cheapest options based on database assets comparisons.", body_style))
        story.append(Spacer(1, 15))

        # Add data table
        if data:
            table_data = [["Item Description", "Best Price", "Cheapest Vendor"]]
            for row in data:
                table_data.append([
                    row.get("item", ""),
                    f"${row.get('cheapest_price', 0.0):.2f}",
                    row.get("vendor", "")
                ])
            
            t = Table(table_data, colWidths=[240, 100, 160])
            t.setStyle(TableStyle([
                ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#0f172a')),
                ('TEXTCOLOR', (0,0), (-1,0), colors.whitesmoke),
                ('ALIGN', (0,0), (-1,-1), 'LEFT'),
                ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
                ('FONTSIZE', (0,0), (-1,0), 10),
                ('BOTTOMPADDING', (0,0), (-1,0), 8),
                ('TOPPADDING', (0,0), (-1,0), 8),
                ('BACKGROUND', (0,1), (-1,-1), colors.HexColor('#f8fafc')),
                ('GRID', (0,0), (-1,-1), 1, colors.HexColor('#e2e8f0')),
                ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
                ('BOTTOMPADDING', (0,1), (-1,-1), 6),
                ('TOPPADDING', (0,1), (-1,-1), 6),
            ]))
            story.append(t)
            
        doc.build(story)
        return True
