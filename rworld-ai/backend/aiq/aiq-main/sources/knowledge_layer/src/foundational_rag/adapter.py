# SPDX-FileCopyrightText: Copyright (c) 2025-2026, NVIDIA CORPORATION & AFFILIATES. All rights reserved.
# SPDX-License-Identifier: Apache-2.0
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
# http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.
"""
Foundational RAG adapter for NVIDIA RAG Blueprint endpoints.

Provides retrieval and ingestion adapters for the hosted RAG Blueprint service.
Requires a deployed RAG Blueprint: https://github.com/NVIDIA-AI-Blueprints/rag

Prerequisites:
    Deploy the NVIDIA RAG Blueprint first:
    https://github.com/NVIDIA-AI-Blueprints/rag/blob/main/docs/deploy-docker-self-hosted.md

Configuration:
    rag_url: Base URL of the RAG server
        - ingestor-server: port 8082 (e.g., "http://ingestor-server:8082/v1")
        - rag-server: port 8081 (e.g., "http://rag-server:8081/v1")
    api_key: Optional API key for authentication
    timeout: Request timeout in seconds (default: 300)

API Reference (ingestor-server :8082):
    - POST /documents - Upload documents (async, returns task_id)
    - GET /documents - List documents in a collection
    - DELETE /documents - Delete documents by name
    - GET /status - Get ingestion task status
    - GET /collections - List all collections
    - POST /collection - Create a new collection
    - DELETE /collections - Delete collections
    - GET /health - Health check

API Reference (rag-server :8081):
    - POST /search - Search/retrieve chunks from collections
    - GET /health - Health check

Repository: https://github.com/NVIDIA-AI-Blueprints/rag
"""

import asyncio
import json
import logging
import os
import re
import threading
import uuid
from datetime import datetime
from functools import partial
from pathlib import Path
from typing import Any

import requests
import urllib3
from requests.adapters import HTTPAdapter
from urllib3.util.retry import Retry

from aiq_agent.knowledge.base import BaseIngestor
from aiq_agent.knowledge.base import BaseRetriever
from aiq_agent.knowledge.base import TTLCleanupMixin
from aiq_agent.knowledge.factory import register_ingestor
from aiq_agent.knowledge.factory import register_retriever
from aiq_agent.knowledge.schema import Chunk
from aiq_agent.knowledge.schema import CollectionInfo
from aiq_agent.knowledge.schema import ContentType
from aiq_agent.knowledge.schema import FileInfo
from aiq_agent.knowledge.schema import FileProgress
from aiq_agent.knowledge.schema import FileStatus
from aiq_agent.knowledge.schema import IngestionJobStatus
from aiq_agent.knowledge.schema import JobState
from aiq_agent.knowledge.schema import RetrievalResult

# Suppress InsecureRequestWarning when verify_ssl=False
urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)

logger = logging.getLogger(__name__)

# Default configuration
# Note: NVIDIA RAG Blueprint has two servers:
#   - Ingestion Server (port 8082): documents, collections, status
#   - Query Server (port 8081): generate, search, retrieval
# @environment_variable RAG_SERVER_URL
# @category Knowledge Layer
# @type str
# @default http://localhost:8081/v1
# @required false
# Base URL for the Foundational RAG query server.
DEFAULT_RAG_URL = os.environ.get("RAG_SERVER_URL", "http://localhost:8081/v1")

# @environment_variable RAG_INGEST_URL
# @category Knowledge Layer
# @type str
# @default http://localhost:8082/v1
# @required false
# Base URL for the Foundational RAG ingestion server.
DEFAULT_INGEST_URL = os.environ.get("RAG_INGEST_URL", "http://localhost:8082/v1")
DEFAULT_TIMEOUT = 300

# Retrieval settings for reranking
VDB_TOP_K_MULTIPLIER = 10
MAX_VDB_TOP_K = 100

# Collection TTL settings (configurable via environment)
COLLECTION_TTL_HOURS = float(os.environ.get("AIQ_COLLECTION_TTL_HOURS", "24"))
TTL_CLEANUP_INTERVAL_SECONDS = int(os.environ.get("AIQ_TTL_CLEANUP_INTERVAL_SECONDS", "3600"))

# Completed jobs are retained for this long so list_files can include failed
# files, then pruned to avoid unbounded memory growth.
JOB_RETENTION_SECONDS = 3600  # 1 hour

# Client-side summary settings (runs parallel to FRAG ingestion)
SUMMARY_MAX_CHARS = 4000
SUMMARY_MAX_PAGES = 2
SUMMARIZABLE_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt", ".md"}


def _create_session(timeout: int = DEFAULT_TIMEOUT, verify_ssl: bool = True) -> requests.Session:
    """Create a requests session with retry logic.

    Args:
        timeout: Request timeout in seconds.
        verify_ssl: Whether to verify SSL certificates. Set False for self-signed certs.
    """
    session = requests.Session()

    # Disable SSL verification if requested (for self-signed certificates)
    session.verify = verify_ssl

    # Configure retry strategy for resilience
    retries = Retry(
        total=3,
        backoff_factor=0.5,
        status_forcelist=[500, 502, 503, 504],
        allowed_methods=["HEAD", "GET", "POST", "PUT", "DELETE", "PATCH"],
    )
    adapter = HTTPAdapter(max_retries=retries)
    session.mount("http://", adapter)
    session.mount("https://", adapter)

    return session


def _extract_text(file_path: str, max_chars: int = SUMMARY_MAX_CHARS) -> str | None:
    """
    Extract text from a file for summary generation.

    Supports PDF, DOCX, PPTX, TXT, and Markdown files.

    Args:
        file_path: Path to the file.
        max_chars: Maximum characters to return.

    Returns:
        Extracted text or None if extraction fails or format is unsupported.
    """
    suffix = Path(file_path).suffix.lower()

    try:
        if suffix == ".pdf":
            try:
                import pypdf
            except ImportError:
                logger.debug("pypdf not installed, skipping PDF text extraction")
                return None
            with open(file_path, "rb") as f:
                reader = pypdf.PdfReader(f)
                text = ""
                for page in reader.pages[:SUMMARY_MAX_PAGES]:
                    text += (page.extract_text() or "") + "\n"
                    if len(text) > max_chars:
                        break
                return text[:max_chars].strip() or None

        elif suffix == ".docx":
            try:
                import docx2txt
            except ImportError:
                logger.debug("docx2txt not installed, skipping DOCX text extraction")
                return None
            text = docx2txt.process(file_path) or ""
            return text[:max_chars].strip() or None

        elif suffix == ".pptx":
            try:
                from pptx import Presentation
            except ImportError:
                logger.debug("python-pptx not installed, skipping PPTX text extraction")
                return None
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text += shape.text_frame.text + "\n"
                    if len(text) > max_chars:
                        break
                if len(text) > max_chars:
                    break
            return text[:max_chars].strip() or None

        elif suffix in (".txt", ".md"):
            with open(file_path, encoding="utf-8", errors="replace") as f:
                text = f.read(max_chars)
            return text.strip() or None

    except Exception as e:
        logger.debug("Text extraction failed for %s: %s", file_path, e)

    return None


def _generate_file_summary(file_path: str, llm=None) -> str | None:
    """
    Generate one-sentence summary from a local file.

    Runs client-side in parallel with FRAG upload. Supports PDF, DOCX,
    PPTX, TXT, and Markdown files.

    Args:
        file_path: Path to the file to summarize.
        llm: LangChain LLM object. Required - no default fallback.

    Returns:
        One-sentence summary or None if unsupported format, no LLM, or generation fails.
    """
    if llm is None:
        return None

    if Path(file_path).suffix.lower() not in SUMMARIZABLE_EXTENSIONS:
        return None

    text = _extract_text(file_path)
    if not text:
        return None

    prompt = f"Summarize in ONE sentence:\n\n{text}"

    try:
        response = llm.invoke(prompt)
        content = response.content if hasattr(response, "content") else str(response)
        return content.strip()
    except Exception as e:
        logger.warning("Summary generation via LLM failed: %s", e)
        return None


@register_retriever("foundational_rag")
class FoundationalRagRetriever(BaseRetriever):
    """
    Retriever adapter that calls hosted NVIDIA RAG Blueprint endpoints.

    This retriever makes HTTP calls to the RAG server's /search endpoint
    and converts responses to the universal Chunk schema.

    Requires a deployed RAG Blueprint server. See:
    https://github.com/NVIDIA-AI-Blueprints/rag/blob/main/docs/deploy-docker-self-hosted.md
    """

    def __init__(self, config: dict[str, Any] | None = None):
        """
        Initialize the RAG server retriever.

        Args:
            config: Configuration dictionary with:
                - rag_url: Base URL of the RAG Query server (port 8081)
                - api_key: Optional API key for authentication
                - timeout: Request timeout in seconds
                - verify_ssl: Whether to verify SSL certificates (default: True)
        """
        super().__init__(config)
        self.rag_url = self.config.get("rag_url", DEFAULT_RAG_URL).rstrip("/")
        # @environment_variable RAG_API_KEY
        # @category API Keys
        # @type str
        # @required false
        # Optional API key for authenticating with the Foundational RAG backend.
        self.api_key = self.config.get("api_key", os.environ.get("RAG_API_KEY"))
        self.timeout = self.config.get("timeout", DEFAULT_TIMEOUT)
        self.verify_ssl = self.config.get("verify_ssl", True)
        self.session = _create_session(self.timeout, self.verify_ssl)

        logger.info(
            f"FoundationalRagRetriever initialized with Query URL: {self.rag_url}, verify_ssl={self.verify_ssl}"
        )

    @property
    def backend_name(self) -> str:
        return "foundational_rag"

    def _get_headers(self) -> dict[str, str]:
        """Get request headers including optional auth."""
        headers = {
            "Content-Type": "application/json",
            "Accept": "application/json",
        }
        if self.api_key:
            headers["Authorization"] = f"Bearer {self.api_key}"
        return headers

    async def retrieve(
        self,
        query: str,
        collection_name: str,
        top_k: int = 10,
        filters: dict[str, Any] | None = None,
    ) -> RetrievalResult:
        """
        Retrieve documents from the RAG server using POST /search.

        The RAG Blueprint search API returns document chunks ranked by relevance.
        Uses reranking for better results.

        Args:
            query: Search query string.
            collection_name: Target collection name.
            top_k: Maximum number of results (maps to reranker_top_k).
            filters: Optional metadata filters (maps to filter_expr).

        Returns:
            RetrievalResult with normalized Chunk objects.
        """
        try:
            endpoint = f"{self.rag_url}/search"

            # Build payload per RAG Blueprint /search API spec
            payload = {
                "query": query,
                "collection_names": [collection_name],  # API expects array
                "reranker_top_k": top_k,  # Final number of results after reranking
                "vdb_top_k": min(top_k * VDB_TOP_K_MULTIPLIER, MAX_VDB_TOP_K),
                "enable_reranker": True,  # Enable reranking for better results
            }

            # Add filter expression if provided
            if filters:
                if isinstance(filters, str):
                    payload["filter_expr"] = filters
                elif "filter_expr" in filters:
                    payload["filter_expr"] = filters["filter_expr"]
                else:
                    # Try to build filter expression from dict
                    # e.g., {"category": "AI"} -> "category == 'AI'"
                    filter_parts = []
                    for k, v in filters.items():
                        if isinstance(v, str):
                            filter_parts.append(f'{k} == "{v}"')
                        else:
                            filter_parts.append(f"{k} == {v}")
                    if filter_parts:
                        payload["filter_expr"] = " and ".join(filter_parts)

            logger.debug(f"Search request to {endpoint}: {payload}")

            # Run blocking HTTP call in thread pool to avoid blocking the event loop
            loop = asyncio.get_event_loop()
            response = await loop.run_in_executor(
                None,
                partial(
                    self.session.post,
                    endpoint,
                    json=payload,
                    headers=self._get_headers(),
                    timeout=self.timeout,
                ),
            )

            response.raise_for_status()
            data = response.json()
            if data is None:
                data = {}

            logger.debug(f"Search response: total_results={data.get('total_results', 0)}")

            # Convert RAG server response to universal schema
            chunks = self._parse_search_response(data, query)

            return RetrievalResult(
                chunks=chunks,
                total_tokens=sum(len(c.content.split()) for c in chunks),
                query=query,
                backend=self.backend_name,
                success=True,
            )

        except requests.exceptions.ConnectionError as e:
            logger.error(f"Cannot connect to RAG server at {self.rag_url}: {e}")
            return RetrievalResult(
                chunks=[],
                total_tokens=0,
                query=query,
                backend=self.backend_name,
                success=False,
                error_message=f"Cannot connect to RAG server: {str(e)[:100]}",
            )

        except requests.exceptions.Timeout as e:
            logger.error(f"RAG server request timed out: {e}")
            return RetrievalResult(
                chunks=[],
                total_tokens=0,
                query=query,
                backend=self.backend_name,
                success=False,
                error_message=f"Request timed out after {self.timeout}s",
            )

        except requests.exceptions.HTTPError as e:
            logger.error(f"RAG server returned error: {e}")
            return RetrievalResult(
                chunks=[],
                total_tokens=0,
                query=query,
                backend=self.backend_name,
                success=False,
                error_message=f"Server error: {str(e)[:100]}",
            )

        except requests.exceptions.RequestException as e:
            logger.error(f"RAG server retrieval failed: {e}")
            return RetrievalResult(
                chunks=[],
                total_tokens=0,
                query=query,
                backend=self.backend_name,
                success=False,
                error_message=f"Request failed: {str(e)[:100]}",
            )

    def _parse_search_response(self, data: dict[str, Any] | None, query: str) -> list[Chunk]:
        """
        Parse RAG server /search response into Chunk objects.

        Response format:
        {
            "total_results": 3,
            "results": [
                {
                    "chunk_id": "...",
                    "document_id": "...",
                    "document_name": "...",
                    "document_type": "text",
                    "content": "...",
                    "page_number": 15,
                    "collection_name": "...",
                    "score": 0.89,
                    "metadata": {...}
                }
            ]
        }
        """
        chunks = []
        if data is None:
            data = {}
        results = data.get("results", [])

        for i, result in enumerate(results):
            chunk = self._normalize_search_result(result, idx=i)
            if chunk:
                chunks.append(chunk)

        return chunks

    def _normalize_search_result(self, result: dict[str, Any], idx: int = 0) -> Chunk | None:
        """
        Convert a single RAG server search result to universal Chunk format.

        Args:
            result: A SourceResult from the /search response.
            idx: Index for fallback chunk ID.

        Returns:
            Normalized Chunk object.
        """
        if not isinstance(result, dict):
            return None

        # Extract fields from the SourceResult schema (use "or {}" so null values become dict)
        document_name_raw = result.get("document_name", "unknown")
        document_type = result.get("document_type", "text").lower()
        content = result.get("content", "")
        score = result.get("score", 0.0)
        collection_name = result.get("collection_name", "")
        metadata = result.get("metadata") or {}
        content_metadata = metadata.get("content_metadata") or {}

        # Strip temp file prefix (tmp + 8 random chars + _) for display; keep raw for delete ops
        display_name = re.sub(r"^tmp.{8}_", "", document_name_raw)

        # Extract page_number from multiple locations (FRAG puts it in metadata)
        page_number = result.get("page_number") or metadata.get("page_number") or content_metadata.get("page_number")

        if page_number in (-1, 0, None):
            page_number = None

        # Generate chunk_id for citation tracking
        doc_base = Path(display_name).stem if display_name != "unknown" else "doc"
        if page_number:
            chunk_id = result.get("chunk_id") or f"{doc_base}_p{page_number}_{idx}"
        else:
            chunk_id = result.get("chunk_id") or f"{doc_base}_{idx}"

        # Determine content type from document_type
        if document_type == "image" or "image" in document_type:
            content_type = ContentType.IMAGE
        elif document_type == "table" or "table" in document_type:
            content_type = ContentType.TABLE
        elif document_type == "chart" or "chart" in document_type:
            content_type = ContentType.CHART
        else:
            content_type = ContentType.TEXT

        # Build display citation
        if page_number and page_number > 0:
            display_citation = f"[{display_name}, p.{page_number}]"
        else:
            display_citation = f"[{display_name}]"

        document_id = result.get("document_id") or document_name_raw

        return Chunk(
            chunk_id=chunk_id,
            content=content if content else "",
            score=float(score),
            file_name=display_name,
            page_number=page_number,
            display_citation=display_citation,
            content_type=content_type,
            metadata={
                "document_id": document_id,
                "document_name_raw": document_name_raw,
                "collection_name": collection_name,
                "page_count": (content_metadata.get("hierarchy") or {}).get("page_count"),
                "language": (content_metadata.get("text_metadata") or {}).get("language"),
                "source_metadata": metadata,
            },
        )

    def normalize(self, raw_result: Any, idx: int = 0) -> Chunk | None:
        """
        Convert RAG server document to universal Chunk format.

        Args:
            raw_result: Raw document from RAG server response.
            idx: Index for generating chunk ID.

        Returns:
            Normalized Chunk object.
        """
        if isinstance(raw_result, dict):
            return self._normalize_search_result(raw_result, idx)

        if isinstance(raw_result, str):
            return Chunk(
                chunk_id=f"rag_{idx}",
                content=raw_result,
                score=1.0,
                file_name="unknown",
                display_citation="[RAG Result]",
                content_type=ContentType.TEXT,
            )

        return None

    async def health_check(self) -> bool:
        """Check if the RAG server is healthy."""
        try:
            loop = asyncio.get_running_loop()
            response = await loop.run_in_executor(
                None,
                lambda: self.session.get(
                    f"{self.rag_url}/health",
                    headers=self._get_headers(),
                    timeout=10,
                ),
            )
            return response.status_code == 200
        except Exception as e:
            logger.warning(f"RAG server health check failed: {e}")
            return False


@register_ingestor("foundational_rag")
class FoundationalRagIngestor(TTLCleanupMixin, BaseIngestor):
    """
    Ingestor adapter that calls hosted NVIDIA RAG Blueprint endpoints.

    This ingestor makes HTTP calls to the RAG server's document management
    and collection management endpoints, enabling remote ingestion and
    file management.

    API Endpoints Used:
        - POST /documents: Upload files for ingestion
        - GET /documents: List files in a collection
        - DELETE /documents: Delete files
        - GET /status: Check ingestion task status
        - GET /collections: List collections
        - POST /collection: Create collection
        - DELETE /collections: Delete collections
    """

    def __init__(self, config: dict[str, Any] | None = None):
        """
        Initialize the RAG server ingestor.

        Args:
            config: Configuration dictionary with:
                - rag_url: Base URL of the RAG Query server (port 8081) - for backwards compat
                - ingest_url: Base URL of the RAG Ingestion server (port 8082)
                - api_key: Optional API key for authentication
                - timeout: Request timeout in seconds
                - chunk_size: Default chunk size for ingestion
                - chunk_overlap: Default chunk overlap
        """
        super().__init__(config)
        # Support both rag_url (legacy) and ingest_url (explicit)
        # If only rag_url is provided, try to derive ingest_url by changing port
        rag_url = self.config.get("rag_url", DEFAULT_RAG_URL).rstrip("/")
        ingest_url = self.config.get("ingest_url", DEFAULT_INGEST_URL).rstrip("/")

        # If ingest_url not explicitly set but rag_url points to 8081, auto-switch to 8082
        if "ingest_url" not in self.config and ":8081" in rag_url:
            ingest_url = rag_url.replace(":8081", ":8082")
        elif "ingest_url" not in self.config and ":8082" in rag_url:
            # User provided 8082 directly as rag_url, use it for ingestion
            ingest_url = rag_url

        self.rag_url = ingest_url  # Ingestor uses ingestion server
        self.api_key = self.config.get("api_key", os.environ.get("RAG_API_KEY"))
        self.timeout = self.config.get("timeout", DEFAULT_TIMEOUT)
        self.chunk_size = self.config.get("chunk_size", 512)
        self.chunk_overlap = self.config.get("chunk_overlap", 150)
        self.generate_summary = self.config.get("generate_summary", False)
        self.summary_llm = self.config.get("summary_llm")  # Resolved LangChain LLM (or None)
        self.verify_ssl = self.config.get("verify_ssl", True)
        self.session = _create_session(self.timeout, self.verify_ssl)

        # Local job tracking (for jobs submitted through this adapter)
        self._jobs: dict[str, IngestionJobStatus] = {}
        self._lock = threading.RLock()

        # Start background TTL cleanup task
        self._start_ttl_cleanup_task(COLLECTION_TTL_HOURS, TTL_CLEANUP_INTERVAL_SECONDS)

        llm_info = "configured" if self.summary_llm else "not configured"
        logger.info(
            "FoundationalRagIngestor initialized: url=%s, generate_summary=%s, summary_llm=%s, verify_ssl=%s",
            self.rag_url,
            self.generate_summary,
            llm_info,
            self.verify_ssl,
        )

    @property
    def backend_name(self) -> str:
        return "foundational_rag"

    def _prune_completed_jobs(self) -> None:
        """Remove completed/failed jobs older than JOB_RETENTION_SECONDS."""
        now = datetime.now()
        with self._lock:
            stale = [
                jid
                for jid, job in self._jobs.items()
                if job.completed_at is not None and (now - job.completed_at).total_seconds() > JOB_RETENTION_SECONDS
            ]
            for jid in stale:
                del self._jobs[jid]
        if stale:
            logger.debug("Pruned %d completed job(s) from tracking", len(stale))

    def _get_headers(self, content_type: str | None = None) -> dict[str, str]:
        """Get request headers including optional auth."""
        headers = {"Accept": "application/json"}
        if content_type:
            headers["Content-Type"] = content_type
        if self.api_key:
            headers["Authorization"] = f"Bearer {self.api_key}"
        return headers

    # =========================================================================
    # Job Submission & Status (Required by BaseIngestor)
    # =========================================================================

    def submit_job(
        self,
        file_paths: list[str],
        collection_name: str,
        config: dict[str, Any] | None = None,
    ) -> str:
        """
        Submit files for ingestion via the RAG server.

        This uploads all files and returns a tracking job ID.
        The RAG server handles ingestion asynchronously.

        Args:
            file_paths: List of local file paths to ingest.
            collection_name: Target collection name.
            config: Optional ingestion configuration (supports original_filenames).

        Returns:
            Job ID for status tracking.
        """
        job_id = str(uuid.uuid4())
        now = datetime.now()
        config = config or {}

        # Original filenames for temp file uploads (avoids tmp prefix)
        original_filenames = config.get("original_filenames", [])

        # Create initial job status using original filenames when available
        file_details = []
        for i, fp in enumerate(file_paths):
            file_name = original_filenames[i] if i < len(original_filenames) else Path(fp).name
            file_details.append(
                FileProgress(
                    file_id=file_name,
                    file_name=file_name,
                    status=FileStatus.UPLOADING,
                )
            )

        job_status = IngestionJobStatus(
            job_id=job_id,
            status=JobState.PENDING,
            submitted_at=now,
            total_files=len(file_paths),
            processed_files=0,
            file_details=file_details,
            collection_name=collection_name,
            backend=self.backend_name,
            metadata={"rag_url": self.rag_url},
        )
        with self._lock:
            self._jobs[job_id] = job_status

        # Upload files in background
        thread = threading.Thread(
            target=self._upload_files_async,
            args=(job_id, file_paths, collection_name, config),
            daemon=True,
        )
        thread.start()

        return job_id

    def _upload_files_async(
        self,
        job_id: str,
        file_paths: list[str],
        collection_name: str,
        config: dict[str, Any] | None = None,
    ):
        """Background thread to upload files to RAG server."""
        from concurrent.futures import ThreadPoolExecutor

        with self._lock:
            job = self._jobs.get(job_id)
            if not job:
                return
            job.status = JobState.PROCESSING
            job.started_at = datetime.now()
        config = config or {}

        # Original filenames for temp file uploads
        original_filenames = config.get("original_filenames", [])

        task_ids = []

        # Track summary futures for parallel generation
        summary_futures: dict[int, tuple[str, Any]] = {}  # index -> (file_name, future)
        executor = None
        if self.generate_summary:
            executor = ThreadPoolExecutor(max_workers=min(len(file_paths), 4))

        # Pre-loop: update statuses and kick off summary generation
        file_handles = []
        files_payload = []
        for i, file_path in enumerate(file_paths):
            file_name = original_filenames[i] if i < len(original_filenames) else Path(file_path).name
            file_path_obj = Path(file_path)

            job.file_details[i].status = FileStatus.INGESTING

            # Start client-side summary generation in parallel (if enabled)
            if self.generate_summary and file_path_obj.suffix.lower() in SUMMARIZABLE_EXTENSIONS:
                logger.info("Starting client-side summary generation")
                future = executor.submit(_generate_file_summary, file_path, self.summary_llm)
                summary_futures[i] = (file_name, future)

            # Open file handle for batch upload
            try:
                fh = open(file_path, "rb")
                file_handles.append(fh)
                files_payload.append(("documents", (file_name, fh, "application/octet-stream")))
            except Exception as e:
                logger.error(f"Failed to open file {file_path}: {e}")
                for fh in file_handles:
                    fh.close()
                raise

        # Single batch upload - all files in one POST request
        # This matches the RAG UI behavior and prevents concurrent ingestion deadlocks
        try:
            data_payload = {
                "collection_name": collection_name,
                "blocking": False,
                "split_options": {
                    "chunk_size": config.get("chunk_size", self.chunk_size),
                    "chunk_overlap": config.get("chunk_overlap", self.chunk_overlap),
                },
                "custom_metadata": [],
                "generate_summary": False,  # Server-side summary disabled, using client-side
            }

            response = self.session.post(
                f"{self.rag_url}/documents",
                files=files_payload,
                data={"data": json.dumps(data_payload)},
                headers=self._get_headers(),
                timeout=self.timeout,
            )

            response.raise_for_status()
            result = response.json()

            # Track the single task ID from the batch upload
            task_id = result.get("task_id", "")
            if task_id:
                task_ids.append(task_id)
                # Map single task_id to all file indices
                job.metadata["task_to_file"] = {task_id: list(range(len(file_paths)))}

            logger.info(f"Uploaded {len(file_paths)} file(s) to RAG server in batch, task_id: {task_id}")

        except Exception as e:
            logger.error(f"Batch upload failed: {e}")
            for i in range(len(file_paths)):
                job.file_details[i].status = FileStatus.FAILED
                job.file_details[i].error_message = str(e)
            job.processed_files = len(file_paths)

        finally:
            for fh in file_handles:
                fh.close()

        # Collect summaries from parallel generation and register them
        if summary_futures:
            from aiq_agent.knowledge import register_summary

            for idx, (file_name, future) in summary_futures.items():
                try:
                    summary = future.result(timeout=30)
                    if summary:
                        register_summary(collection_name, file_name, summary)
                        logger.info(f"  Summary generated ({len(summary)} chars)")
                except Exception as e:
                    logger.debug(f"Summary generation failed for {file_name}: {e}")

        # Clean up executor
        if executor:
            executor.shutdown(wait=False)

        # Store task IDs for status polling
        job.metadata["task_ids"] = task_ids

        # Check if all uploads failed immediately
        failed_count = sum(1 for fd in job.file_details if fd.status == FileStatus.FAILED)
        if failed_count == job.total_files:
            job.status = JobState.FAILED
            job.error_message = "File upload failed"
            job.completed_at = datetime.now()
        # Otherwise keep as PROCESSING - get_job_status will poll FRAG and update

    def get_job_status(self, job_id: str) -> IngestionJobStatus:
        """
        Get the status of an ingestion job.

        This checks local job tracking and optionally polls the
        RAG server for task status updates.

        Args:
            job_id: Job ID from submit_job.

        Returns:
            Current job status.
        """
        self._prune_completed_jobs()

        with self._lock:
            job = self._jobs.get(job_id)
        if not job:
            # Return a failed status for unknown jobs
            return IngestionJobStatus(
                job_id=job_id,
                status=JobState.FAILED,
                submitted_at=datetime.now(),
                collection_name="unknown",
                backend=self.backend_name,
                error_message=f"Job {job_id} not found",
            )

        # If job is still processing, check RAG server task status
        if job.status == JobState.PROCESSING:
            task_ids = job.metadata.get("task_ids", [])
            task_to_file = job.metadata.get("task_to_file", {})

            for task_id in task_ids:
                # Get file indices for this task (list for batch, int for legacy single-file)
                file_idx_raw = task_to_file.get(task_id)
                if isinstance(file_idx_raw, list):
                    file_indices = file_idx_raw
                elif file_idx_raw is not None:
                    file_indices = [file_idx_raw]
                else:
                    file_indices = []

                try:
                    response = self.session.get(
                        f"{self.rag_url}/status",
                        params={"task_id": task_id},
                        headers=self._get_headers("application/json"),
                        timeout=30,
                    )
                    if response.status_code == 200:
                        status_data = response.json()
                        state = status_data.get("state", "").lower()
                        logger.debug(f"FRAG task {task_id[:8]} state={state}")

                        if state in ("success", "completed", "finished"):
                            result = status_data.get("result", {})
                            failed_docs = result.get("failed_documents", [])
                            failed_names: dict[str, str] = {}
                            for fdoc in failed_docs:
                                fname = fdoc.get("document_name", "")
                                ferr = fdoc.get("error_message", "Ingestion failed")
                                if fname:
                                    failed_names[fname] = ferr

                            if failed_names:
                                logger.warning(
                                    f"FRAG task {task_id[:8]} completed with "
                                    f"{len(failed_names)} failed file(s): {list(failed_names.keys())}"
                                )

                            if file_indices:
                                for idx in file_indices:
                                    if (
                                        idx < len(job.file_details)
                                        and job.file_details[idx].status == FileStatus.INGESTING
                                    ):
                                        fd = job.file_details[idx]
                                        # Check if this file appears in failed_documents
                                        matched_error = next(
                                            (err for fname, err in failed_names.items() if fname == fd.file_name),
                                            None,
                                        )
                                        if matched_error:
                                            fd.status = FileStatus.FAILED
                                            fd.error_message = matched_error
                                        else:
                                            fd.status = FileStatus.SUCCESS
                                succeeded = sum(
                                    1
                                    for i in file_indices
                                    if i < len(job.file_details) and job.file_details[i].status == FileStatus.SUCCESS
                                )
                                failed = sum(
                                    1
                                    for i in file_indices
                                    if i < len(job.file_details) and job.file_details[i].status == FileStatus.FAILED
                                )
                                logger.debug(f"Batch ingestion done: {succeeded} succeeded, {failed} failed")
                            else:
                                # Fallback: match by document name
                                docs = result.get("documents", [])
                                for doc in docs:
                                    doc_name = doc.get("document_name", "")
                                    for fd in job.file_details:
                                        if doc_name == fd.file_name and fd.status == FileStatus.INGESTING:
                                            fd.status = FileStatus.SUCCESS
                                            logger.debug("File ingestion succeeded")
                                for fname, ferr in failed_names.items():
                                    for fd in job.file_details:
                                        if fname == fd.file_name and fd.status == FileStatus.INGESTING:
                                            fd.status = FileStatus.FAILED
                                            fd.error_message = ferr
                                            logger.warning(f"File ingestion failed: {ferr}")
                        elif state == "failed":
                            result = status_data.get("result", {})
                            # Try multiple fields for error message
                            # FRAG puts detailed errors in result.message
                            error_msg = (
                                status_data.get("error")
                                or status_data.get("error_message")
                                or status_data.get("message")
                                or result.get("message")  # FRAG uses this field
                                or result.get("error")
                                or result.get("error_message")
                                or "Unknown error"
                            )
                            failed_docs = result.get("failed_documents", [])
                            logger.warning(f"FRAG task {task_id[:8]} failed: {error_msg}")

                            # Mark all files in this batch task as failed
                            if file_indices:
                                for idx in file_indices:
                                    if (
                                        idx < len(job.file_details)
                                        and job.file_details[idx].status == FileStatus.INGESTING
                                    ):
                                        job.file_details[idx].status = FileStatus.FAILED
                                        job.file_details[idx].error_message = error_msg
                            elif failed_docs:
                                # Fallback: use failed_docs list
                                for fdoc in failed_docs:
                                    doc_name = fdoc.get("document_name", "")
                                    doc_error = fdoc.get("error_message", error_msg)
                                    for fd in job.file_details:
                                        if doc_name == fd.file_name and fd.status == FileStatus.INGESTING:
                                            fd.status = FileStatus.FAILED
                                            fd.error_message = doc_error
                                            logger.warning(f"File ingestion failed: {doc_error}")
                        elif state in ("pending", "started", "processing", "running"):
                            # Check per-document status for incremental progress
                            # The RAG server provides nv_ingest_status.document_wise_status
                            # which shows individual file completion before the batch finishes
                            nv_status = status_data.get("nv_ingest_status", {})
                            doc_statuses = nv_status.get("document_wise_status", {})
                            if doc_statuses:
                                for fd in job.file_details:
                                    if fd.status == FileStatus.INGESTING:
                                        doc_state = doc_statuses.get(fd.file_name, "")
                                        if doc_state == "completed":
                                            fd.status = FileStatus.SUCCESS
                                            logger.debug(f"File {fd.file_name} completed (batch still processing)")
                    else:
                        logger.warning(f"FRAG status check returned {response.status_code} for task {task_id[:8]}")
                except Exception as e:
                    logger.warning(f"Failed to check FRAG task {task_id[:8]} status: {e}")

            # After checking all tasks, update job completion status
            success_count = sum(1 for fd in job.file_details if fd.status == FileStatus.SUCCESS)
            failed_count = sum(1 for fd in job.file_details if fd.status == FileStatus.FAILED)
            total_done = success_count + failed_count

            if total_done == job.total_files:
                # All files have terminal status
                job.processed_files = total_done
                if failed_count == job.total_files:
                    job.status = JobState.FAILED
                    job.error_message = "File ingestion failed"
                else:
                    job.status = JobState.COMPLETED
                job.completed_at = datetime.now()

                logger.info(f"Job {job_id[:8]} completed: {success_count} succeeded, {failed_count} failed")

        return job

    # =========================================================================
    # Collection Management
    # =========================================================================

    def create_collection(
        self,
        name: str,
        description: str | None = None,
        metadata: dict[str, Any] | None = None,
    ) -> CollectionInfo:
        """
        Create a new collection via the RAG server.

        Args:
            name: Collection name (typically session_id).
            description: Human-readable description.
            metadata: Additional configuration (embedding_dimension, etc.).

        Returns:
            CollectionInfo with created collection details.
        """
        try:
            # Use the new /collection endpoint
            payload = {
                "collection_name": name,
                "embedding_dimension": metadata.get("embedding_dimension", 2048) if metadata else 2048,
                "metadata_schema": metadata.get("metadata_schema", []) if metadata else [],
            }

            # Optionally specify VDB endpoint if provided
            if metadata and "vdb_endpoint" in metadata:
                payload["vdb_endpoint"] = metadata["vdb_endpoint"]

            response = self.session.post(
                f"{self.rag_url}/collection",
                json=payload,
                headers=self._get_headers("application/json"),
                timeout=self.timeout,
            )
            response.raise_for_status()
            result = response.json()

            logger.info(f"Created collection: {name}")

            return CollectionInfo(
                name=name,
                description=description or result.get("message", ""),
                file_count=0,
                chunk_count=0,
                created_at=datetime.now(),
                backend=self.backend_name,
                metadata={"rag_url": self.rag_url, **(metadata or {})},
            )

        except requests.exceptions.HTTPError as e:
            if e.response.status_code == 422:
                # Collection might already exist, try to get it
                existing = self.get_collection(name)
                if existing:
                    return existing
            logger.error(f"Failed to create collection {name}: {e}")
            raise

    def delete_collection(self, name: str) -> bool:
        """
        Delete a collection via the RAG server.

        Args:
            name: Collection name to delete.

        Returns:
            True if deleted successfully.
        """
        try:
            response = self.session.delete(
                f"{self.rag_url}/collections",
                json=[name],  # API expects array of collection names
                headers=self._get_headers("application/json"),
                timeout=self.timeout,
            )
            response.raise_for_status()
            result = response.json()

            successful = result.get("successful", [])
            if name in successful:
                # Clear summaries from centralized registry
                from aiq_agent.knowledge import clear_collection_summaries

                clear_collection_summaries(name)

                logger.info(f"Deleted collection: {name}")
                return True

            failed = result.get("failed", [])
            for f in failed:
                if f.get("collection_name") == name:
                    logger.warning(f"Failed to delete collection {name}: {f.get('error_message')}")

            return False

        except Exception as e:
            logger.error(f"Failed to delete collection {name}: {e}")
            return False

    def list_collections(self) -> list[CollectionInfo]:
        """
        List all collections from the RAG server.

        Returns:
            List of CollectionInfo objects with timestamps from FRAG's collection_info.
        """
        try:
            response = self.session.get(
                f"{self.rag_url}/collections",
                headers=self._get_headers("application/json"),
                timeout=self.timeout,
            )
            response.raise_for_status()
            data = response.json()

            collections = []
            for coll in data.get("collections", []):
                # Extract collection_info which contains timestamps and metadata
                coll_info = coll.get("collection_info", {})

                # Parse timestamps from FRAG response
                created_at = None
                updated_at = None

                if coll_info.get("date_created"):
                    try:
                        created_at = datetime.fromisoformat(coll_info["date_created"].replace("Z", "+00:00"))
                    except (ValueError, TypeError):
                        pass

                # Map last_indexed to updated_at for 24hr clearance logic
                # Priority: last_indexed (ingestion time) > last_updated > date_created
                timestamp_str = (
                    coll_info.get("last_indexed") or coll_info.get("last_updated") or coll_info.get("date_created")
                )
                if timestamp_str:
                    try:
                        updated_at = datetime.fromisoformat(timestamp_str.replace("Z", "+00:00"))
                    except (ValueError, TypeError):
                        pass

                collections.append(
                    CollectionInfo(
                        name=coll.get("collection_name", "unknown"),
                        description=coll_info.get("description") or None,
                        file_count=coll_info.get("number_of_files", 0),
                        chunk_count=coll.get("num_entities", 0),
                        created_at=created_at,
                        updated_at=updated_at,
                        backend=self.backend_name,
                        metadata={
                            "rag_url": self.rag_url,
                            "metadata_schema": coll.get("metadata_schema", []),
                            "last_indexed": coll_info.get("last_indexed"),
                            "ingestion_status": coll_info.get("ingestion_status"),
                            "doc_type_counts": coll_info.get("doc_type_counts"),
                        },
                    )
                )

            return collections

        except Exception as e:
            logger.error(f"Failed to list collections: {e}")
            return []

    def get_collection(self, name: str) -> CollectionInfo | None:
        """
        Get metadata for a specific collection.

        Args:
            name: Collection name.

        Returns:
            CollectionInfo if found, None otherwise.
        """
        collections = self.list_collections()
        for coll in collections:
            if coll.name == name:
                return coll
        logger.debug(f"Collection '{name}' not found in {len(collections)} available collections")
        return None

    # =========================================================================
    # File Management
    # =========================================================================

    def upload_file(
        self,
        file_path: str,
        collection_name: str,
        metadata: dict[str, Any] | None = None,
    ) -> FileInfo:
        """
        Upload a file to a collection via the RAG server.

        This submits the file for async ingestion and returns immediately.
        Use get_file_status to poll for completion.

        If generate_summary is enabled, a client-side summary is generated
        in parallel with the upload using local text extraction and NIM.

        Args:
            file_path: Local path to the file.
            collection_name: Target collection (session_id).
            metadata: Optional file metadata.

        Returns:
            FileInfo with upload status and optional summary.
        """
        from concurrent.futures import ThreadPoolExecutor

        file_path_obj = Path(file_path)
        file_id = str(uuid.uuid4())
        now = datetime.now()

        # Initial file info
        file_info = FileInfo(
            file_id=file_id,
            file_name=file_path_obj.name,
            collection_name=collection_name,
            status=FileStatus.UPLOADING,
            file_size=file_path_obj.stat().st_size if file_path_obj.exists() else None,
            chunk_count=0,
            uploaded_at=now,
            metadata=metadata or {},
        )

        # Start client-side summary generation in background (parallel with upload)
        summary_future = None
        executor = None
        if self.generate_summary and file_path_obj.suffix.lower() in SUMMARIZABLE_EXTENSIONS:
            logger.info("Starting client-side summary generation for %s", file_path_obj.name)
            executor = ThreadPoolExecutor(max_workers=1)
            summary_future = executor.submit(_generate_file_summary, file_path, self.summary_llm)

        try:
            # Build the data payload
            data_payload = {
                "collection_name": collection_name,
                "blocking": False,
                "split_options": {
                    "chunk_size": metadata.get("chunk_size", self.chunk_size) if metadata else self.chunk_size,
                    "chunk_overlap": metadata.get("chunk_overlap", self.chunk_overlap)
                    if metadata
                    else self.chunk_overlap,
                },
                "custom_metadata": metadata.get("custom_metadata", []) if metadata else [],
                "generate_summary": metadata.get("generate_summary", False) if metadata else False,
            }

            with open(file_path, "rb") as f:
                files = {"documents": (file_path_obj.name, f, "application/octet-stream")}
                form_data = {"data": json.dumps(data_payload)}

                response = self.session.post(
                    f"{self.rag_url}/documents",
                    files=files,
                    data=form_data,
                    headers=self._get_headers(),
                    timeout=self.timeout,
                )

            response.raise_for_status()
            result = response.json()

            # Store task_id for status tracking
            task_id = result.get("task_id", "")
            file_info.status = FileStatus.INGESTING
            file_info.metadata["task_id"] = task_id
            file_info.metadata["message"] = result.get("message", "")

            logger.info(f"Uploaded file, task_id: {task_id}")

        except Exception as e:
            file_info.status = FileStatus.FAILED
            file_info.error_message = str(e)
            logger.error(f"Failed to upload file {file_path}: {e}")

        # Wait for client-side summary (non-blocking during upload)
        if summary_future:
            try:
                summary = summary_future.result(timeout=15)
                file_info.metadata["summary"] = summary

                # Register in centralized summary registry (backend-agnostic)
                if summary:
                    from aiq_agent.knowledge import register_summary

                    register_summary(collection_name, file_info.file_name, summary)
                    logger.info(f"  Summary generated ({len(summary)} chars)")

            except TimeoutError:
                logger.warning("Summary generation timed out for %s", file_path_obj.name)
                file_info.metadata["summary"] = None
            except Exception as e:
                logger.warning("Summary generation failed for %s: %s", file_path_obj.name, e)
                file_info.metadata["summary"] = None

        # Clean up executor
        if executor:
            executor.shutdown(wait=False)

        return file_info

    def delete_file(self, file_id: str, collection_name: str) -> bool:
        """
        Delete a file from a collection via the RAG server.

        Args:
            file_id: File identifier (typically filename/document_name).
            collection_name: Collection containing the file.

        Returns:
            True if deleted successfully.
        """
        try:
            response = self.session.delete(
                f"{self.rag_url}/documents",
                params={"collection_name": collection_name},
                json=[file_id],  # API expects array of document names
                headers=self._get_headers("application/json"),
                timeout=self.timeout,
            )
            response.raise_for_status()
            result = response.json()

            # Check if this file was in the deleted documents
            docs = result.get("documents", [])
            remaining_names = [d.get("document_name", "") for d in docs]

            # If the file is no longer in the list, it was deleted
            if file_id not in remaining_names:
                # Remove from centralized summary registry
                from aiq_agent.knowledge import unregister_summary

                unregister_summary(collection_name, file_id)

                logger.info(f"Deleted file {file_id} from collection {collection_name}")
                return True

            logger.warning(f"File {file_id} still exists after delete request")
            return False

        except Exception as e:
            logger.error(f"Failed to delete file {file_id}: {e}")
            return False

    def delete_files(
        self,
        file_ids: list[str],
        collection_name: str,
    ) -> dict[str, Any]:
        """
        Delete multiple files from a collection (batch delete).

        Args:
            file_ids: List of file IDs (document_names) to delete.
            collection_name: Collection containing the files.

        Returns:
            Dict with 'successful', 'failed', 'total_deleted'.
        """
        try:
            response = self.session.delete(
                f"{self.rag_url}/documents",
                params={"collection_name": collection_name},
                json=file_ids,
                headers=self._get_headers("application/json"),
                timeout=self.timeout,
            )
            response.raise_for_status()
            result = response.json()

            # Trust FRAG's 200 response as confirmation of successful delete
            # Remove from centralized summary registry
            from aiq_agent.knowledge import unregister_summary

            for file_id in file_ids:
                unregister_summary(collection_name, file_id)

            # Also purge failed files from local job tracking so they don't
            # reappear via list_files (failed files only exist in _jobs, not
            # on the FRAG server, so the FRAG delete above is a no-op for them)
            deleted_set = set(file_ids)
            with self._lock:
                for job in self._jobs.values():
                    if job.collection_name != collection_name:
                        continue
                    job.file_details = [fd for fd in job.file_details if fd.file_name not in deleted_set]

            logger.info(f"Deleted {len(file_ids)} files from {collection_name}")
            return {
                "message": result.get("message", f"Deleted {len(file_ids)} files"),
                "successful": file_ids,
                "failed": [],
                "total_deleted": len(file_ids),
            }
        except Exception as e:
            logger.error(f"Failed to delete files from {collection_name}: {e}")
            return {
                "message": f"Delete failed: {e}",
                "successful": [],
                "failed": [{"file_id": fid, "error": str(e)} for fid in file_ids],
                "total_deleted": 0,
            }

    def list_files(self, collection_name: str) -> list[FileInfo]:
        """
        List all files in a collection via the RAG server.

        Args:
            collection_name: Collection to list files from.

        Returns:
            List of FileInfo objects.
        """
        try:
            response = self.session.get(
                f"{self.rag_url}/documents",
                params={"collection_name": collection_name},
                headers=self._get_headers("application/json"),
                timeout=self.timeout,
            )
            response.raise_for_status()
            data = response.json()

            files = []
            for doc in data.get("documents", []):
                doc_name = doc.get("document_name", "unknown")
                doc_metadata = doc.get("metadata", {})
                doc_info = doc.get("document_info", {})

                # Extract fields from document_info
                file_size = doc_info.get("file_size") if doc_info else None
                chunk_count = doc_info.get("total_elements", 0) if doc_info else 0
                ingested_at = None
                if doc_info and doc_info.get("date_created"):
                    try:
                        ingested_at = datetime.fromisoformat(doc_info["date_created"])
                    except (ValueError, TypeError):
                        pass

                # Merge document_info into metadata for UI access
                enriched_metadata = {**doc_metadata}
                if doc_info:
                    enriched_metadata["document_type"] = doc_info.get("document_type")
                    enriched_metadata["doc_type_counts"] = doc_info.get("doc_type_counts")
                    enriched_metadata["raw_text_elements_size"] = doc_info.get("raw_text_elements_size")

                files.append(
                    FileInfo(
                        file_id=doc_name,  # RAG server uses document_name as ID
                        file_name=doc_name,
                        collection_name=collection_name,
                        status=FileStatus.SUCCESS,
                        file_size=file_size,
                        chunk_count=chunk_count,
                        uploaded_at=ingested_at,
                        ingested_at=ingested_at,
                        metadata=enriched_metadata,
                    )
                )

            # Include failed files from local job tracking.
            # The FRAG server only returns successfully ingested documents, so
            # files that failed ingestion would otherwise vanish from the UI.
            existing_names = {f.file_name for f in files}
            with self._lock:
                jobs_snapshot = list(self._jobs.values())
            for job in jobs_snapshot:
                if job.collection_name != collection_name:
                    continue
                for fd in job.file_details:
                    if fd.status == FileStatus.FAILED and fd.file_name not in existing_names:
                        files.append(
                            FileInfo(
                                file_id=fd.file_name,
                                file_name=fd.file_name,
                                collection_name=collection_name,
                                status=FileStatus.FAILED,
                                error_message=fd.error_message,
                            )
                        )
                        existing_names.add(fd.file_name)

            return files

        except Exception as e:
            logger.error(f"Failed to list files in collection {collection_name}: {e}")
            return []

    def get_file_status(self, file_id: str, collection_name: str) -> FileInfo | None:
        """
        Get the status of a specific file using the RAG server's /status endpoint.

        The RAG server uses Celery task states:
        - PENDING: Task is waiting to be processed
        - STARTED: Task is currently being processed
        - SUCCESS/FINISHED: Task completed successfully
        - FAILURE/FAILED: Task failed

        Args:
            file_id: File identifier (task_id from upload response).
            collection_name: Collection containing the file.

        Returns:
            FileInfo with current status, None if task not found.
        """
        try:
            response = self.session.get(
                f"{self.rag_url}/status",
                params={"task_id": file_id},
                headers=self._get_headers("application/json"),
                timeout=30,
            )

            if response.status_code == 200:
                status_data = response.json()
                # RAG server returns uppercase Celery states
                state = status_data.get("state", "UNKNOWN").upper()
                result = status_data.get("result", {})

                logger.debug(f"Task {file_id} status: state={state}, result={result}")

                # Map Celery states to FileStatus
                # Terminal states: SUCCESS, FINISHED, FAILURE, FAILED
                # Non-terminal states: PENDING, STARTED
                if state in ("SUCCESS", "FINISHED"):
                    file_status = FileStatus.SUCCESS
                elif state in ("FAILURE", "FAILED"):
                    file_status = FileStatus.FAILED
                elif state == "STARTED":
                    file_status = FileStatus.INGESTING
                elif state == "PENDING":
                    file_status = FileStatus.INGESTING  # Still waiting/processing
                else:
                    file_status = FileStatus.INGESTING  # Unknown = assume processing

                # Extract document info from result (available when finished)
                docs = result.get("documents", [])
                failed_docs = result.get("failed_documents", [])
                total_docs = result.get("total_documents", 0)
                message = result.get("message", "")

                # Determine filename and error message
                doc_name = file_id  # Default to task_id
                error_msg = None

                if docs:
                    # Success case - get document name from result
                    doc_name = docs[0].get("document_name", file_id)
                elif failed_docs:
                    # Failure case - get error info
                    doc_name = failed_docs[0].get("document_name", file_id)
                    error_msg = failed_docs[0].get("error_message", "Ingestion failed")
                    file_status = FileStatus.FAILED
                elif state in ("FAILURE", "FAILED"):
                    # Generic failure without document details
                    error_msg = message or "Ingestion failed"

                return FileInfo(
                    file_id=file_id,
                    file_name=doc_name,
                    collection_name=collection_name,
                    status=file_status,
                    chunk_count=total_docs,
                    error_message=error_msg,
                    metadata={
                        "raw_state": state,
                        "message": message,
                        "result": result,
                    },
                )

            elif response.status_code == 404:
                logger.debug(f"Task {file_id} not found (404)")
                return None

        except requests.exceptions.RequestException as e:
            logger.debug(f"Task status request failed: {e}")
        except Exception as e:
            logger.debug(f"Task status lookup failed: {e}")

        # Fall back to listing files and finding by name (for completed files)
        files = self.list_files(collection_name)
        for f in files:
            if file_id in (f.file_id, f.file_name):
                return f

        return None

    async def health_check(self) -> bool:
        """Check if the RAG server is healthy."""
        try:
            loop = asyncio.get_running_loop()
            response = await loop.run_in_executor(
                None,
                lambda: self.session.get(
                    f"{self.rag_url}/health",
                    headers=self._get_headers(),
                    timeout=10,
                ),
            )
            if response.status_code == 200:
                data = response.json()
                # Check if message indicates healthy status
                return "up" in data.get("message", "").lower() or response.status_code == 200
            return False
        except Exception as e:
            logger.warning(f"RAG server health check failed: {e}")
            return False
