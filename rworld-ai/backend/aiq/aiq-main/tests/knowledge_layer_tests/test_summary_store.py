# SPDX-FileCopyrightText: Copyright (c) 2025-2026, NVIDIA CORPORATION & AFFILIATES. All rights reserved.
# SPDX-License-Identifier: Apache-2.0
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
# http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

"""Unit tests for document summarization storage.

Tests cover:
- SummaryStore SQLAlchemy-based storage (SQLite)
- Factory functions (register_summary, get_available_documents, etc.)
- AvailableDocument model
- URL normalization for database connections
"""

import tempfile
from pathlib import Path

import pytest

from aiq_agent.knowledge.schema import AvailableDocument
from aiq_agent.knowledge.summary_store import SummaryStore
from aiq_agent.knowledge.summary_store import _normalize_db_url

# =============================================================================
# URL Normalization Tests
# =============================================================================


class TestNormalizeDbUrl:
    """Tests for the _normalize_db_url helper function."""

    def test_sqlite_url_async_mode(self):
        """Test SQLite URL normalization for async mode."""
        url = "sqlite:///./summaries.db"
        result = _normalize_db_url(url, async_mode=True)
        assert result == "sqlite+aiosqlite:///./summaries.db"

    def test_sqlite_url_sync_mode(self):
        """Test SQLite URL normalization for sync mode."""
        url = "sqlite:///./summaries.db"
        result = _normalize_db_url(url, async_mode=False)
        assert result == "sqlite:///./summaries.db"

    def test_sqlite_already_async(self):
        """Test SQLite URL that already has async driver."""
        url = "sqlite+aiosqlite:///./summaries.db"
        result = _normalize_db_url(url, async_mode=True)
        assert result == "sqlite+aiosqlite:///./summaries.db"

    def test_postgresql_url_async_mode(self):
        """Test PostgreSQL URL normalization for async mode."""
        url = "postgresql://user:pass@localhost:5432/db"
        result = _normalize_db_url(url, async_mode=True)
        assert result == "postgresql+psycopg://user:pass@localhost:5432/db"

    def test_postgresql_url_sync_mode(self):
        """Test PostgreSQL URL normalization for sync mode."""
        url = "postgresql://user:pass@localhost:5432/db"
        result = _normalize_db_url(url, async_mode=False)
        assert result == "postgresql+psycopg://user:pass@localhost:5432/db"

    def test_postgres_shorthand_url(self):
        """Test postgres:// shorthand URL normalization."""
        url = "postgres://user:pass@localhost:5432/db"
        result = _normalize_db_url(url, async_mode=True)
        assert result == "postgresql+psycopg://user:pass@localhost:5432/db"

    def test_postgresql_with_existing_driver(self):
        """Test PostgreSQL URL with existing driver gets normalized."""
        url = "postgresql+asyncpg://user:pass@localhost:5432/db"
        result = _normalize_db_url(url, async_mode=True)
        assert "psycopg" in result

    def test_unknown_url_passthrough(self):
        """Test unknown database URLs pass through unchanged."""
        url = "mysql://user:pass@localhost/db"
        result = _normalize_db_url(url, async_mode=True)
        assert result == url


# =============================================================================
# AvailableDocument Model Tests
# =============================================================================


class TestAvailableDocument:
    """Tests for the AvailableDocument Pydantic model."""

    def test_create_with_summary(self):
        """Test creating AvailableDocument with a summary."""
        doc = AvailableDocument(file_name="test.pdf", summary="A test document.")
        assert doc.file_name == "test.pdf"
        assert doc.summary == "A test document."

    def test_create_without_summary(self):
        """Test creating AvailableDocument without a summary."""
        doc = AvailableDocument(file_name="test.pdf")
        assert doc.file_name == "test.pdf"
        assert doc.summary is None

    def test_model_dump(self):
        """Test model serialization to dict."""
        doc = AvailableDocument(file_name="report.pdf", summary="Financial report.")
        data = doc.model_dump()
        assert data == {"file_name": "report.pdf", "summary": "Financial report."}

    def test_model_dump_without_summary(self):
        """Test model serialization without summary."""
        doc = AvailableDocument(file_name="report.pdf")
        data = doc.model_dump()
        assert data == {"file_name": "report.pdf", "summary": None}

    def test_model_validate(self):
        """Test model creation from dict."""
        data = {"file_name": "doc.pdf", "summary": "Test summary"}
        doc = AvailableDocument.model_validate(data)
        assert doc.file_name == "doc.pdf"
        assert doc.summary == "Test summary"


# =============================================================================
# SummaryStore Tests
# =============================================================================


class TestSummaryStore:
    """Tests for the SummaryStore SQLAlchemy-based storage."""

    @pytest.fixture
    def temp_db(self):
        """Create a temporary SQLite database for testing."""
        with tempfile.TemporaryDirectory() as tmpdir:
            db_path = Path(tmpdir) / "test_summaries.db"
            db_url = f"sqlite:///{db_path}"
            yield db_url

    @pytest.fixture
    def store(self, temp_db):
        """Create a SummaryStore instance with temp database."""
        return SummaryStore(temp_db)

    def test_store_initialization(self, temp_db):
        """Test SummaryStore initializes correctly."""
        store = SummaryStore(temp_db)
        assert store.db_url == temp_db

    def test_register_summary(self, store):
        """Test registering a document summary."""
        store.register("test_collection", "doc1.pdf", "This is a test summary.")
        docs = store.get_all("test_collection")
        assert len(docs) == 1
        assert docs[0].file_name == "doc1.pdf"
        assert docs[0].summary == "This is a test summary."

    def test_register_multiple_summaries(self, store):
        """Test registering multiple document summaries."""
        store.register("collection1", "doc1.pdf", "Summary 1")
        store.register("collection1", "doc2.pdf", "Summary 2")
        store.register("collection1", "doc3.pdf", "Summary 3")

        docs = store.get_all("collection1")
        assert len(docs) == 3
        filenames = {doc.file_name for doc in docs}
        assert filenames == {"doc1.pdf", "doc2.pdf", "doc3.pdf"}

    def test_register_updates_existing(self, store):
        """Test registering a summary for existing file updates it."""
        store.register("collection", "doc.pdf", "Original summary")
        store.register("collection", "doc.pdf", "Updated summary")

        docs = store.get_all("collection")
        assert len(docs) == 1
        assert docs[0].summary == "Updated summary"

    def test_get_all_empty_collection(self, store):
        """Test getting documents from empty collection returns empty list."""
        docs = store.get_all("nonexistent_collection")
        assert docs == []

    def test_get_all_different_collections(self, store):
        """Test documents are isolated by collection."""
        store.register("collection_a", "doc1.pdf", "Summary A1")
        store.register("collection_a", "doc2.pdf", "Summary A2")
        store.register("collection_b", "doc3.pdf", "Summary B1")

        docs_a = store.get_all("collection_a")
        docs_b = store.get_all("collection_b")

        assert len(docs_a) == 2
        assert len(docs_b) == 1
        assert docs_b[0].file_name == "doc3.pdf"

    def test_unregister_summary(self, store):
        """Test unregistering a document summary."""
        store.register("collection", "doc1.pdf", "Summary 1")
        store.register("collection", "doc2.pdf", "Summary 2")

        store.unregister("collection", "doc1.pdf")

        docs = store.get_all("collection")
        assert len(docs) == 1
        assert docs[0].file_name == "doc2.pdf"

    def test_unregister_nonexistent(self, store):
        """Test unregistering nonexistent document doesn't raise error."""
        store.unregister("collection", "nonexistent.pdf")  # Should not raise

    def test_clear_collection(self, store):
        """Test clearing all summaries in a collection."""
        store.register("collection", "doc1.pdf", "Summary 1")
        store.register("collection", "doc2.pdf", "Summary 2")
        store.register("other_collection", "doc3.pdf", "Summary 3")

        store.clear_collection("collection")

        assert store.get_all("collection") == []
        assert len(store.get_all("other_collection")) == 1

    def test_clear_all(self, store):
        """Test clearing all summaries across all collections."""
        store.register("collection1", "doc1.pdf", "Summary 1")
        store.register("collection2", "doc2.pdf", "Summary 2")

        store.clear_all()

        assert store.get_all("collection1") == []
        assert store.get_all("collection2") == []

    @pytest.mark.asyncio
    async def test_get_all_async(self, store):
        """Test async retrieval of document summaries."""
        store.register("async_collection", "doc1.pdf", "Async summary 1")
        store.register("async_collection", "doc2.pdf", "Async summary 2")

        docs = await store.get_all_async("async_collection")

        assert len(docs) == 2
        filenames = {doc.file_name for doc in docs}
        assert filenames == {"doc1.pdf", "doc2.pdf"}

    @pytest.mark.asyncio
    async def test_get_all_async_empty(self, store):
        """Test async retrieval from empty collection."""
        docs = await store.get_all_async("empty_collection")
        assert docs == []


# =============================================================================
# Factory Function Tests
# =============================================================================


class TestFactoryFunctions:
    """Tests for the factory module's summary registry functions."""

    @pytest.fixture(autouse=True)
    def reset_summary_store(self):
        """Reset the global summary store before each test."""
        from aiq_agent.knowledge import factory

        factory._summary_store = None
        yield
        factory._summary_store = None

    @pytest.fixture
    def temp_db_url(self):
        """Create a temporary SQLite database URL."""
        with tempfile.TemporaryDirectory() as tmpdir:
            db_path = Path(tmpdir) / "factory_test.db"
            yield f"sqlite:///{db_path}"

    def test_configure_summary_db(self, temp_db_url):
        """Test configuring the summary database."""
        from aiq_agent.knowledge import configure_summary_db

        configure_summary_db(temp_db_url)

        from aiq_agent.knowledge import factory

        assert factory._summary_store is not None
        assert factory._summary_store.db_url == temp_db_url

    def test_register_summary(self, temp_db_url):
        """Test registering a summary via factory function."""
        from aiq_agent.knowledge import configure_summary_db
        from aiq_agent.knowledge import get_available_documents
        from aiq_agent.knowledge.factory import register_summary

        configure_summary_db(temp_db_url)
        register_summary("test_collection", "test.pdf", "Test summary")

        docs = get_available_documents("test_collection")
        assert len(docs) == 1
        assert docs[0].file_name == "test.pdf"
        assert docs[0].summary == "Test summary"

    def test_register_summary_none_skipped(self, temp_db_url):
        """Test that None summaries are not registered."""
        from aiq_agent.knowledge import configure_summary_db
        from aiq_agent.knowledge import get_available_documents
        from aiq_agent.knowledge.factory import register_summary

        configure_summary_db(temp_db_url)
        register_summary("collection", "doc.pdf", None)

        docs = get_available_documents("collection")
        assert len(docs) == 0

    def test_register_summary_empty_string_skipped(self, temp_db_url):
        """Test that empty string summaries are not registered."""
        from aiq_agent.knowledge import configure_summary_db
        from aiq_agent.knowledge import get_available_documents
        from aiq_agent.knowledge.factory import register_summary

        configure_summary_db(temp_db_url)
        register_summary("collection", "doc.pdf", "")

        docs = get_available_documents("collection")
        assert len(docs) == 0

    def test_get_available_documents(self, temp_db_url):
        """Test getting available documents via factory function."""
        from aiq_agent.knowledge import configure_summary_db
        from aiq_agent.knowledge import get_available_documents
        from aiq_agent.knowledge.factory import register_summary

        configure_summary_db(temp_db_url)
        register_summary("my_collection", "file1.pdf", "Summary for file 1")
        register_summary("my_collection", "file2.pdf", "Summary for file 2")

        docs = get_available_documents("my_collection")
        assert len(docs) == 2
        assert all(isinstance(doc, AvailableDocument) for doc in docs)

    @pytest.mark.asyncio
    async def test_get_available_documents_async(self, temp_db_url):
        """Test async retrieval of available documents."""
        from aiq_agent.knowledge import configure_summary_db
        from aiq_agent.knowledge.factory import get_available_documents_async
        from aiq_agent.knowledge.factory import register_summary

        configure_summary_db(temp_db_url)
        register_summary("async_test", "doc.pdf", "Async test summary")

        docs = await get_available_documents_async("async_test")
        assert len(docs) == 1
        assert docs[0].file_name == "doc.pdf"

    def test_unregister_summary(self, temp_db_url):
        """Test unregistering a summary via factory function."""
        from aiq_agent.knowledge import configure_summary_db
        from aiq_agent.knowledge import get_available_documents
        from aiq_agent.knowledge.factory import register_summary
        from aiq_agent.knowledge.factory import unregister_summary

        configure_summary_db(temp_db_url)
        register_summary("collection", "doc1.pdf", "Summary 1")
        register_summary("collection", "doc2.pdf", "Summary 2")

        unregister_summary("collection", "doc1.pdf")

        docs = get_available_documents("collection")
        assert len(docs) == 1
        assert docs[0].file_name == "doc2.pdf"

    def test_clear_collection_summaries(self, temp_db_url):
        """Test clearing collection summaries via factory function."""
        from aiq_agent.knowledge import configure_summary_db
        from aiq_agent.knowledge import get_available_documents
        from aiq_agent.knowledge.factory import clear_collection_summaries
        from aiq_agent.knowledge.factory import register_summary

        configure_summary_db(temp_db_url)
        register_summary("coll1", "doc1.pdf", "Summary 1")
        register_summary("coll2", "doc2.pdf", "Summary 2")

        clear_collection_summaries("coll1")

        assert get_available_documents("coll1") == []
        assert len(get_available_documents("coll2")) == 1

    def test_clear_all_summaries(self, temp_db_url):
        """Test clearing all summaries via factory function."""
        from aiq_agent.knowledge import configure_summary_db
        from aiq_agent.knowledge import get_available_documents
        from aiq_agent.knowledge.factory import clear_all_summaries
        from aiq_agent.knowledge.factory import register_summary

        configure_summary_db(temp_db_url)
        register_summary("coll1", "doc1.pdf", "Summary 1")
        register_summary("coll2", "doc2.pdf", "Summary 2")

        clear_all_summaries()

        assert get_available_documents("coll1") == []
        assert get_available_documents("coll2") == []

    def test_lazy_initialization(self):
        """Test that summary store is lazily initialized with default DB."""
        from aiq_agent.knowledge import get_available_documents

        # Don't call configure_summary_db - should auto-initialize
        docs = get_available_documents("test_collection")
        assert docs == []  # Empty but works

        from aiq_agent.knowledge import factory

        assert factory._summary_store is not None


# =============================================================================
# Integration Tests
# =============================================================================


class TestSummaryIntegration:
    """Integration tests for the summary storage workflow."""

    @pytest.fixture
    def temp_db_url(self):
        """Create a temporary SQLite database URL."""
        with tempfile.TemporaryDirectory() as tmpdir:
            db_path = Path(tmpdir) / "integration_test.db"
            yield f"sqlite:///{db_path}"

    @pytest.fixture(autouse=True)
    def reset_store(self):
        """Reset global store before each test."""
        from aiq_agent.knowledge import factory

        factory._summary_store = None
        yield
        factory._summary_store = None

    def test_full_workflow(self, temp_db_url):
        """Test complete summary storage workflow."""
        from aiq_agent.knowledge import configure_summary_db
        from aiq_agent.knowledge import get_available_documents
        from aiq_agent.knowledge.factory import clear_collection_summaries
        from aiq_agent.knowledge.factory import register_summary
        from aiq_agent.knowledge.factory import unregister_summary

        # 1. Configure database
        configure_summary_db(temp_db_url)

        # 2. Simulate document ingestion with summaries
        collection = "session_12345"
        register_summary(collection, "report.pdf", "Annual financial report for 2024.")
        register_summary(collection, "data.csv", "Sales data spreadsheet.")
        register_summary(collection, "notes.txt", "Meeting notes from Q4 review.")

        # 3. Retrieve for agent context
        docs = get_available_documents(collection)
        assert len(docs) == 3

        # 4. User deletes a file
        unregister_summary(collection, "notes.txt")
        docs = get_available_documents(collection)
        assert len(docs) == 2

        # 5. Session ends, cleanup
        clear_collection_summaries(collection)
        docs = get_available_documents(collection)
        assert len(docs) == 0

    def test_multiple_collections_isolation(self, temp_db_url):
        """Test that different sessions/collections are isolated."""
        from aiq_agent.knowledge import configure_summary_db
        from aiq_agent.knowledge import get_available_documents
        from aiq_agent.knowledge.factory import register_summary

        configure_summary_db(temp_db_url)

        # Two different user sessions
        session_a = "user_alice_session_001"
        session_b = "user_bob_session_002"

        register_summary(session_a, "alice_doc.pdf", "Alice's document")
        register_summary(session_b, "bob_doc.pdf", "Bob's document")
        register_summary(session_b, "bob_doc2.pdf", "Bob's second document")

        # Each user only sees their own documents
        alice_docs = get_available_documents(session_a)
        bob_docs = get_available_documents(session_b)

        assert len(alice_docs) == 1
        assert alice_docs[0].file_name == "alice_doc.pdf"

        assert len(bob_docs) == 2
        assert {d.file_name for d in bob_docs} == {"bob_doc.pdf", "bob_doc2.pdf"}

    @pytest.mark.asyncio
    async def test_mixed_sync_async_operations(self, temp_db_url):
        """Test mixing sync registration with async retrieval."""
        from aiq_agent.knowledge import configure_summary_db
        from aiq_agent.knowledge.factory import get_available_documents_async
        from aiq_agent.knowledge.factory import register_summary

        configure_summary_db(temp_db_url)

        # Sync registration (as backend adapters do)
        register_summary("mixed_test", "doc1.pdf", "Document 1 summary")
        register_summary("mixed_test", "doc2.pdf", "Document 2 summary")

        # Async retrieval (as agent code does)
        docs = await get_available_documents_async("mixed_test")

        assert len(docs) == 2
        assert all(doc.summary is not None for doc in docs)


# =============================================================================
# Text Extraction Tests (Foundational RAG multi-format support)
# =============================================================================


class TestExtractText:
    """Tests for the _extract_text helper that supports PDF, DOCX, PPTX, TXT, and MD."""

    def test_txt_extraction(self):
        """Test text extraction from a .txt file."""
        from knowledge_layer.foundational_rag.adapter import _extract_text

        with tempfile.NamedTemporaryFile(suffix=".txt", mode="w", delete=False, encoding="utf-8") as f:
            f.write("Hello world. This is a test document for summarization.")
            f.flush()
            result = _extract_text(f.name)

        assert result is not None
        assert "Hello world" in result

    def test_md_extraction(self):
        """Test text extraction from a .md file."""
        from knowledge_layer.foundational_rag.adapter import _extract_text

        with tempfile.NamedTemporaryFile(suffix=".md", mode="w", delete=False, encoding="utf-8") as f:
            f.write("# Heading\n\nSome markdown content with **bold** text.")
            f.flush()
            result = _extract_text(f.name)

        assert result is not None
        assert "Heading" in result
        assert "bold" in result

    def test_txt_max_chars_truncation(self):
        """Test that text extraction respects max_chars limit."""
        from knowledge_layer.foundational_rag.adapter import _extract_text

        with tempfile.NamedTemporaryFile(suffix=".txt", mode="w", delete=False, encoding="utf-8") as f:
            f.write("A" * 10000)
            f.flush()
            result = _extract_text(f.name, max_chars=100)

        assert result is not None
        assert len(result) <= 100

    def test_empty_txt_returns_none(self):
        """Test that empty files return None."""
        from knowledge_layer.foundational_rag.adapter import _extract_text

        with tempfile.NamedTemporaryFile(suffix=".txt", mode="w", delete=False, encoding="utf-8") as f:
            f.write("")
            f.flush()
            result = _extract_text(f.name)

        assert result is None

    def test_unsupported_extension_returns_none(self):
        """Test that unsupported file extensions return None."""
        from knowledge_layer.foundational_rag.adapter import _extract_text

        with tempfile.NamedTemporaryFile(suffix=".csv", mode="w", delete=False, encoding="utf-8") as f:
            f.write("col1,col2\nval1,val2")
            f.flush()
            result = _extract_text(f.name)

        assert result is None

    def test_nonexistent_file_returns_none(self):
        """Test that a nonexistent file returns None gracefully."""
        from knowledge_layer.foundational_rag.adapter import _extract_text

        result = _extract_text("/nonexistent/path/file.txt")
        assert result is None

    def test_pdf_extraction(self):
        """Test PDF text extraction using pypdf."""
        from knowledge_layer.foundational_rag.adapter import _extract_text

        pytest.importorskip("pypdf")
        pdf_dir = Path(__file__).parent / "data"
        pdf_files = list(pdf_dir.glob("*.pdf"))
        if not pdf_files:
            pytest.skip("No test PDFs in data/ directory")

        result = _extract_text(str(pdf_files[0]))
        assert result is not None
        assert len(result) > 0

    def test_docx_extraction(self):
        """Test DOCX text extraction using docx2txt."""
        from knowledge_layer.foundational_rag.adapter import _extract_text

        pytest.importorskip("docx2txt")

        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
            tmp_path = f.name

        try:
            from docx import Document

            doc = Document()
            doc.add_paragraph("Test paragraph for DOCX extraction.")
            doc.save(tmp_path)

            result = _extract_text(tmp_path)
            assert result is not None
            assert "Test paragraph" in result
        except ImportError:
            pytest.skip("python-docx not installed for DOCX creation")

    def test_pptx_extraction(self):
        """Test PPTX text extraction using python-pptx."""
        from knowledge_layer.foundational_rag.adapter import _extract_text

        pytest.importorskip("pptx")
        from pptx import Presentation

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            tmp_path = f.name

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Test Slide Title"
        slide.placeholders[1].text = "Slide body content for testing."
        prs.save(tmp_path)

        result = _extract_text(tmp_path)
        assert result is not None
        assert "Test Slide Title" in result


# =============================================================================
# File Summary Generation Tests
# =============================================================================


class TestGenerateFileSummary:
    """Tests for _generate_file_summary with multi-format support."""

    def test_no_llm_returns_none(self):
        """Test that None LLM returns None."""
        from knowledge_layer.foundational_rag.adapter import _generate_file_summary

        with tempfile.NamedTemporaryFile(suffix=".txt", mode="w", delete=False, encoding="utf-8") as f:
            f.write("Some content")
            f.flush()
            result = _generate_file_summary(f.name, llm=None)

        assert result is None

    def test_unsupported_format_returns_none(self):
        """Test that unsupported file formats return None even with LLM."""
        from unittest.mock import MagicMock

        from knowledge_layer.foundational_rag.adapter import _generate_file_summary

        mock_llm = MagicMock()
        with tempfile.NamedTemporaryFile(suffix=".csv", mode="w", delete=False, encoding="utf-8") as f:
            f.write("data")
            f.flush()
            result = _generate_file_summary(f.name, llm=mock_llm)

        assert result is None
        mock_llm.invoke.assert_not_called()

    def test_txt_summary_calls_llm(self):
        """Test that a .txt file triggers LLM summarization."""
        from unittest.mock import MagicMock

        from knowledge_layer.foundational_rag.adapter import _generate_file_summary

        mock_llm = MagicMock()
        mock_llm.invoke.return_value = MagicMock(content="A one-sentence summary.")

        with tempfile.NamedTemporaryFile(suffix=".txt", mode="w", delete=False, encoding="utf-8") as f:
            f.write("This is a long document about artificial intelligence and machine learning.")
            f.flush()
            result = _generate_file_summary(f.name, llm=mock_llm)

        assert result == "A one-sentence summary."
        mock_llm.invoke.assert_called_once()

    def test_md_summary_calls_llm(self):
        """Test that a .md file triggers LLM summarization."""
        from unittest.mock import MagicMock

        from knowledge_layer.foundational_rag.adapter import _generate_file_summary

        mock_llm = MagicMock()
        mock_llm.invoke.return_value = MagicMock(content="Markdown summary.")

        with tempfile.NamedTemporaryFile(suffix=".md", mode="w", delete=False, encoding="utf-8") as f:
            f.write("# Research Notes\n\nKey findings about neural networks.")
            f.flush()
            result = _generate_file_summary(f.name, llm=mock_llm)

        assert result == "Markdown summary."
        mock_llm.invoke.assert_called_once()

    def test_llm_failure_returns_none(self):
        """Test graceful handling when LLM raises an exception."""
        from unittest.mock import MagicMock

        from knowledge_layer.foundational_rag.adapter import _generate_file_summary

        mock_llm = MagicMock()
        mock_llm.invoke.side_effect = RuntimeError("LLM API error")

        with tempfile.NamedTemporaryFile(suffix=".txt", mode="w", delete=False, encoding="utf-8") as f:
            f.write("Some document content.")
            f.flush()
            result = _generate_file_summary(f.name, llm=mock_llm)

        assert result is None

    def test_summarizable_extensions_constant(self):
        """Test that SUMMARIZABLE_EXTENSIONS contains expected formats."""
        from knowledge_layer.foundational_rag.adapter import SUMMARIZABLE_EXTENSIONS

        assert ".pdf" in SUMMARIZABLE_EXTENSIONS
        assert ".docx" in SUMMARIZABLE_EXTENSIONS
        assert ".pptx" in SUMMARIZABLE_EXTENSIONS
        assert ".txt" in SUMMARIZABLE_EXTENSIONS
        assert ".md" in SUMMARIZABLE_EXTENSIONS
        assert ".csv" not in SUMMARIZABLE_EXTENSIONS
